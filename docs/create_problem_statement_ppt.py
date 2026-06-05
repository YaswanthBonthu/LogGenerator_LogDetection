"""Generate Problem Statement PowerPoint for ThreatScope project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = r"c:\Users\Dell\OneDrive\Desktop\New folder\docs\ThreatScope_Problem_Statement.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(8, 11, 20)
CYAN = RGBColor(0, 212, 255)
WHITE = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
VIOLET = RGBColor(124, 58, 237)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CYAN
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(16)


def add_bullet_slide(title, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)
        p.level = 0

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12), Inches(0.4))
        fb.text_frame.paragraphs[0].text = footer
        fb.text_frame.paragraphs[0].font.size = Pt(12)
        fb.text_frame.paragraphs[0].font.color.rgb = VIOLET


add_title_slide(
    "ThreatScope",
    "AI-Powered Security Log Evaluation & Threat Intelligence Platform",
)

add_bullet_slide("Problem Statement", [
    "Organizations generate massive volumes of logs across authentication, web servers, applications, firewalls, and network devices.",
    "Security teams struggle to manually correlate events, identify attack patterns, and prioritize real threats in time.",
    "Known attacks repeat across systems, but most SIEM tools lack per-system memory of previously seen malicious behavior.",
    "CVE data exists in NVD but is rarely linked automatically to live log evidence and attack context.",
    "Alert fatigue: too many generic alerts with template remediation — not actionable for executives or engineers.",
])

add_bullet_slide("Current Pain Points", [
    "Slow incident response — analysts spend hours parsing raw logs manually",
    "Unknown threats (zero-day behavior) missed by signature-only rule engines",
    "No unified view: detection + vulnerability + reasoning + remediation in one flow",
    "Disconnected tools: log generators, evaluators, and dashboards are separate silos",
    "Non-technical stakeholders cannot understand severity without plain-English summaries",
])

add_bullet_slide("Proposed Solution — Two-App Architecture", [
    "App 1: SecureCorp Dummy Website — simulates auth, web, app, firewall & network with continuous log generation",
    "App 2: ThreatScope Log Evaluator — separate application consuming logs via live feed API",
    "Decoupled design: generator exposes /logs/recent; evaluator pulls via /analyze/live",
    "Per-system threat memory (SQLite) — recognizes previously seen attacks on each host",
    "React dashboards for both applications with real-time monitoring",
])

add_bullet_slide("ThreatScope Pipeline", [
    "1. Log Ingestion + Parsing (JSON, NDJSON, CSV, Syslog)",
    "2. Rule Engine (fast signatures) + ML Anomaly Detector (Isolation Forest)",
    "3. CVE Correlation — software fingerprinting, NVD API + SQLite cache, EPSS scores",
    "4. Reasoning Layer — Gemini 2.0 Flash: attack stage, exploitable CVEs, remediation, blast radius",
    "5. Alert Generator — CVSS + attack stage severity, evidence with line numbers, plain-English summary",
])

add_bullet_slide("Key Features Delivered", [
    "Live log feed from dummy server → evaluator (no manual file upload required)",
    "Fast mode: sub-second rule-based scan; full mode: CVE + AI reasoning in background",
    "Threat memory per host (web-01, db-01, etc.) with KNOWN threat badges",
    "MITRE-aligned detection: brute force, SQLi, XSS, port scan, C2, privilege escalation",
    "Evidence-linked alerts with ordered remediation from AI reasoning (not templates)",
])

add_bullet_slide("Technology Stack", [
    "Backend: Python, FastAPI, scikit-learn, SQLite (CVE cache + threat memory)",
    "AI: Google Gemini 2.0 Flash with deterministic fallback reasoning",
    "Frontend: React + Vite + Chart.js (ThreatScope & SecureCorp dashboards)",
    "Data sources: NVD REST API, FIRST EPSS API, internal continuous log generator",
    "Deployment: Local dev — ports 8100 (logs), 8000 (evaluator API), 5173/5180 (UIs)",
])

add_bullet_slide("Expected Outcomes & Impact", [
    "Reduce mean-time-to-detect (MTTD) via automated rule + ML hybrid detection",
    "Prioritize patching using CVE + EPSS correlation tied to observed software versions",
    "Executive-ready incident summaries without sacrificing technical evidence depth",
    "Repeat-attack recognition lowers noise and highlights persistent adversaries",
    "Demonstrates end-to-end SOC workflow in a self-contained demo environment",
])

add_bullet_slide("Demo Workflow", [
    "1. Start SecureCorp dummy website — logs generate continuously in background",
    "2. User interacts (login, API calls) — more auth/web/app events are recorded",
    "3. Open ThreatScope → Connect Live Feed",
    "4. Quick scan shows threats in seconds; full CVE + AI analysis loads in background",
    "5. Review threats, memory, CVEs, and line-linked evidence in dashboard tabs",
], footer="ThreatScope v2.0 — SecureCorp + Evaluator")

prs.save(OUT)
print(f"Created: {OUT}")
